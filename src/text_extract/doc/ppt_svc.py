import pptx
from Util import split_sentences_from_paragraph

from text_extract.doc.abc_doc_extract import AbstractDocExtractSvc


class PptSvc(AbstractDocExtractSvc):
    def __init__(self):
        super().__init__()

    def extract_from_doc(self, path: str):
        prs = pptx.Presentation(path)
        sentence_list = list()
        for i, slide in enumerate(prs.slides, start=1):
            for j, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, "text"):
                    sentence_list.extend(split_sentences_from_paragraph(shape.text))

        # Remove duplicates
        sentence_list = list(dict.fromkeys(sentence_list))

        return sentence_list


ppt_extract_svc = PptSvc()
